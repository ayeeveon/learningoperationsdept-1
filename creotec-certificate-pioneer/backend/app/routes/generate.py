from flask import Blueprint, request, jsonify, send_file
from flask_cors import cross_origin
from pptx import Presentation
import zipfile
import os
import json
import io
from copy import deepcopy

bp = Blueprint('generate', __name__, url_prefix='/generate')
# Template and output folders relative to this file
# Determine project base and correct template/output paths
BASE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
TEMPLATE_DIR = os.path.join(BASE_DIR, 'uploads', 'templates')
OUTPUT_DIR = os.path.join(BASE_DIR, 'static', 'generated')

os.makedirs(OUTPUT_DIR, exist_ok=True)

# Add helper to fill placeholders using entire text frame
def fill_slide(slide, data_row):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            # Combine runs text to handle placeholders split across runs
            full_text = ''.join(run.text for run in paragraph.runs)
            replaced = full_text
            for key, value in data_row.items():
                ph = f"{{{key}}}"
                replaced = replaced.replace(ph, str(value))
            if replaced != full_text:
                # Clear existing run texts
                for run in paragraph.runs:
                    run.text = ""
                # Set replacement text in first run
                paragraph.runs[0].text = replaced

@bp.route('/certificates', methods=['POST', 'OPTIONS'])
@cross_origin()
def generate_certificates():
    data = request.get_json()
    template_type = data.get('template', 'ojt')
    rows = data.get('rows', [])

    if not rows:
        return jsonify({"error": "No data provided"}), 400

    # Determine template file
    if template_type in ['ojt', 'immersion']:
        tpl_filename = f"{template_type}_default.pptx"
    else:
        tpl_filename = f"{template_type}.pptx"
    template_path = os.path.join(TEMPLATE_DIR, tpl_filename)
    if not os.path.exists(template_path):
        return jsonify({"error": f"Template '{tpl_filename}' not found"}), 404

    # Load template presentation
    prs = Presentation(template_path)

    # Prepare source slide and clone elements
    source_slide = prs.slides[0]
    # Deepcopy all shape XML elements from source slide to preserve template content
    original_elements = [deepcopy(shape.element) for shape in source_slide.shapes]
    # Fill first slide with first row
    # Fill first slide
    fill_slide(source_slide, rows[0])

    # Create and fill slides for remaining rows by cloning source slide elements
    for row in rows[1:]:
        new_slide = prs.slides.add_slide(source_slide.slide_layout)
        # Remove default layout shapes
        for shp in list(new_slide.shapes):
            new_slide.shapes._spTree.remove(shp.element)
        # Append cloned shapes from source slide
        for el in original_elements:
            new_slide.shapes._spTree.append(deepcopy(el))
        # Fill placeholders on the cloned slide
        fill_slide(new_slide, row)

    # Save combined PPTX
    output_name = f"certificates_{template_type}.pptx"
    output_path = os.path.join(OUTPUT_DIR, output_name)
    prs.save(output_path)
    return jsonify({"message": "Certificates generated", "files": [output_name]})

@bp.route('/files/<filename>', methods=['GET'])
@cross_origin()
def get_generated_file(filename):
    file_path = os.path.join(OUTPUT_DIR, filename)
    if not os.path.exists(file_path):
        return jsonify({"error": "File not found"}), 404
    return send_file(file_path, as_attachment=True)

@bp.route('/preview', methods=['POST', 'OPTIONS'])
@cross_origin()
def preview_certificate():
    data = request.get_json()
    template_type = data.get('template', 'ojt')
    rows = data.get('rows', [])
    if not rows:
        return jsonify({"error": "No data to preview"}), 400
    # Determine template filename
    if template_type in ['ojt', 'immersion']:
        filename = f"{template_type}_default.pptx"
    else:
        filename = f"{template_type}.pptx"
    template_path = os.path.join(TEMPLATE_DIR, filename)
    if not os.path.exists(template_path):
        return jsonify({"error": f"Template '{filename}' not found"}), 404
    # Build HTML preview with enhanced styling
    html_parts = [
        "<!DOCTYPE html>",
        "<html>",
        "<head>",
        "<meta charset='utf-8'>",
        "<meta name='viewport' content='width=device-width, initial-scale=1.0'>",
        "<title>Certificate Preview</title>",
        "<style>",
        "body { font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; background: #2d3748; margin: 0; padding: 20px; min-height: 100vh; }",
        ".container { max-width: 900px; margin: 0 auto; }",
        "h2 { text-align: center; color: white; font-size: 2.5rem; margin-bottom: 30px; text-shadow: 2px 2px 4px rgba(0,0,0,0.3); }",
        ".slide-preview { background: #4a5568; border-radius: 12px; padding: 25px; margin-bottom: 25px; box-shadow: 0 8px 32px rgba(0,0,0,0.2); border: 1px solid #718096; transition: transform 0.2s ease; }",
        ".slide-preview:hover { transform: translateY(-2px); box-shadow: 0 12px 40px rgba(0,0,0,0.3); }",
        ".slide-preview h4 { color: white; font-size: 1.3rem; margin: 0 0 15px 0; padding-bottom: 8px; border-bottom: 2px solid #718096; }",
        ".slide-preview p { margin: 10px 0; font-size: 1.1rem; line-height: 1.6; color: #e2e8f0; padding: 8px 12px; background: #2d3748; border-radius: 6px; border-left: 4px solid #a361ef; }",
        ".certificate-text { font-weight: 500; }",
        "@media (max-width: 768px) { .container { padding: 10px; } .slide-preview { padding: 20px; } h2 { font-size: 2rem; } }",
        "</style>",
        "</head>",
        "<body>",
        "<div class='container'>",
        "<h2>Certificate Preview</h2>"
    ]
    for idx, row in enumerate(rows):
        # Render slide text after placeholder replacement
        prs_row = Presentation(template_path)
        slide = prs_row.slides[0]
        fill_slide(slide, row)
        html_parts.append("<div class='slide-preview certificate-text'>")
        html_parts.append(f"<h4>Certificate {idx+1}</h4>")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = ''.join(run.text for run in paragraph.runs)
                if text.strip():
                    html_parts.append(f"<p>{text}</p>")
        html_parts.append("</div>")
    html_parts.append("</div>")
    html_parts.append("</body></html>")
    return "\n".join(html_parts), 200, {"Content-Type": "text/html"}

